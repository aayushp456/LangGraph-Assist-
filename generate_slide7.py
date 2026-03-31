from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE

# ── Palette ──────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
DARK_RED   = RGBColor(0x8B, 0x1A, 0x1A)
BLUE       = RGBColor(0x2E, 0x75, 0xB6)
CRAG_RED   = RGBColor(0xB5, 0x3A, 0x1A)
GRAY_BOX   = RGBColor(0x75, 0x75, 0x75)
STAT_BG    = RGBColor(0xE2, 0xEF, 0xDA)
STAT_GREEN = RGBColor(0x37, 0x5E, 0x0F)
STAT_BORD  = RGBColor(0x70, 0xAD, 0x47)
TEXT_GRAY  = RGBColor(0x50, 0x50, 0x50)
WHITE_DIM  = RGBColor(0xD8, 0xD8, 0xD8)
ARROW_COL  = RGBColor(0x80, 0x80, 0x80)

# ── Presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shp.line.color.rgb = line_rgb
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def set_tf(shp, lines, anchor=MSO_ANCHOR.MIDDLE):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for item in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = item.get("align", PP_ALIGN.LEFT)
        for seg in item.get("runs", []):
            r = p.add_run()
            r.text = seg["text"]
            r.font.size = Pt(seg.get("size", 11))
            r.font.bold = seg.get("bold", False)
            r.font.italic = seg.get("italic", False)
            r.font.color.rgb = seg.get("color", BLACK)


def textbox(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    set_tf(tb, lines, anchor)
    return tb


# ── Title ─────────────────────────────────────────────────────────────────────
textbox(slide, 0.3, 0.13, 12.7, 0.65, [
    {"runs": [{"text": "The RAG Pipeline — Why I Went Beyond Basic Vector Search",
               "size": 22, "bold": True, "color": BLACK}],
     "align": PP_ALIGN.LEFT}
])

# Red separator line
add_rect(slide, 0.3, 0.85, 12.7, 0.04, DARK_RED)

# ── LEFT PANEL ────────────────────────────────────────────────────────────────
# Opening quote
textbox(slide, 0.3, 0.98, 5.15, 1.05, [
    {"runs": [{"text": (
        "\u201cNaive cosine similarity returned articles that were topically "
        "related but factually wrong \u2014 a ticket about a Python SDK timeout "
        "got a JavaScript rate-limiting article.\u201d"
    ), "size": 10.5, "italic": True, "color": TEXT_GRAY}],
     "align": PP_ALIGN.LEFT}
], anchor=MSO_ANCHOR.TOP)

# Sub-heading
textbox(slide, 0.3, 2.08, 5.15, 0.3, [
    {"runs": [{"text": "So I engineered a corrective multi-stage pipeline:",
               "size": 11, "bold": True, "color": BLACK}],
     "align": PP_ALIGN.LEFT}
])

# 3 bullets
bullets = [
    ("Query Expansion",
     " \u2014 Gemini generates 2\u20133 rephrased variants before hitting "
     "Pinecone, widening recall without losing precision."),
    ("CRAG Evaluator",
     " \u2014 Grades each retrieved doc. Scores < 0.4 trigger a live Tavily "
     "web search fallback \u2014 no hallucination from bad context."),
    ("CrossEncoder + MMR Reranking",
     " \u2014 Joint (query, doc) scoring + diversity so top-5 aren\u2019t "
     "all variations of the same article."),
]

by = 2.44
for title, desc in bullets:
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(by), Inches(5.15), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "\u25b8 " + title
    r1.font.size = Pt(10.5)
    r1.font.bold = True
    r1.font.color.rgb = BLUE
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(10)
    r2.font.color.rgb = TEXT_GRAY
    by += 0.70

# Stat box
stat_shp = add_rect(slide, 0.3, 4.60, 5.15, 0.65, STAT_BG, STAT_BORD)
set_tf(stat_shp, [
    {"runs": [{"text": "  \U0001f4c8  40\u201350% improvement in retrieval precision "
                       "vs. basic vector search  ",
               "size": 11.5, "bold": True, "color": STAT_GREEN}],
     "align": PP_ALIGN.CENTER}
])

# ── RIGHT PANEL: 8-stage pipeline ────────────────────────────────────────────
stages = [
    (GRAY_BOX,
     "Stage 1 \u2014 Semantic Cache Check",
     "cosine \u2265 0.92 \u2192 serve in ~0.1s  \u00b7  skip all stages below   \u2193 MISS"),
    (BLUE,
     "Stage 2 \u2014 Query Expansion  (Gemini)",
     "generates 2\u20133 query variants to widen recall"),
    (BLUE,
     "Stage 3 \u2014 Pinecone Vector Search",
     "top-k = 15 candidates per variant  \u00b7  cosine similarity"),
    (BLUE,
     "Stage 4 \u2014 Hybrid Keyword Boost",
     "score = 0.7 \u00d7 vector  +  0.3 \u00d7 BM25 keyword"),
    (BLUE,
     "Stage 5 \u2014 Metadata Filter",
     "category match \u00d71.3  \u00b7  recency boost \u00d71.15  \u00b7  score threshold"),
    (CRAG_RED,
     "\u2605  Stage 6 \u2014 CRAG Evaluator  (Gemini)",
     "CORRECT \u2192 proceed   |   AMBIGUOUS \u2192 refine + re-retrieve   |   "
     "INCORRECT \u2192 Tavily web search"),
    (BLUE,
     "Stage 7 \u2014 CrossEncoder + MMR Reranker",
     "joint (query, doc) relevance scoring  +  diversity selection  \u03bb = 0.7"),
    (GRAY_BOX,
     "Stage 8 \u2014 Cache Store  +  Return Top-5",
     "TTL = 1 h  \u00b7  max 500 entries  \u00b7  feeds solution generator"),
]

rx, rw = 5.65, 7.40
box_h  = 0.545
gap    = 0.185
sy     = 0.98

for i, (col, name, desc) in enumerate(stages):
    yp = sy + i * (box_h + gap)

    # Box
    shp = add_rect(slide, rx, yp, rw, box_h, col)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = "  " + name + "   "
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(9)
    r2.font.color.rgb = WHITE_DIM

    # Connector arrow
    if i < len(stages) - 1:
        ax = Inches(rx + rw / 2)
        ay1 = Inches(yp + box_h)
        ay2 = Inches(yp + box_h + gap)
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT, ax, ay1, ax, ay2
        )
        conn.line.color.rgb = ARROW_COL
        conn.line.width = Pt(1.5)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/aayushpatel/src/Support-Agent/slide7_rag_pipeline.pptx"
prs.save(out)
print(f"Saved: {out}")
